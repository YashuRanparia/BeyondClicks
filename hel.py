from pptx import Presentation
from pptx.util import Inches
from PyQt5.QtCore import Qt, pyqtSlot
from PyQt5.QtGui import QPixmap
import aspose.slides
from PyQt5.QtWidgets import (QGraphicsScene, QGraphicsView, QPushButton,
                             QVBoxLayout, QWidget)


class PresentationScreen(QWidget):
    def __init__(self, pptx_file_path):
        super().__init__()
        self.presentation = Presentation(pptx_file_path)
        self.current_slide_index = 0
        self.setup_ui()

    def setup_ui(self):
        layout = QVBoxLayout()

        # Display the current slide's content
        self.slide_view = QGraphicsView()
        layout.addWidget(self.slide_view)

        # Next Slide button
        next_button = QPushButton("Next Slide")
        next_button.clicked.connect(self.show_next_slide)
        layout.addWidget(next_button)

        self.setLayout(layout)
        self.show_current_slide()

    def show_current_slide(self):
        current_slide = self.presentation.slides[self.current_slide_index]
        slide_image = self.get_slide_image(current_slide)
        self.slide_view.setScene(slide_image)

    def get_slide_image(self, slide):
        # Calculate slide dimensions (adjust as needed)
        slide_width = Inches(6.5)
        slide_height = Inches(4.5)

        # Render the slide as an image
        image_stream = slide._element.get_or_add_image_part().blob
        pixmap = QPixmap()
        pixmap.loadFromData(image_stream)
        pixmap = pixmap.scaled(slide_width, slide_height, aspectMode=Qt.KeepAspectRatio, transformMode=Qt.SmoothTransformation)
        scene = QGraphicsScene()
        scene.addPixmap(pixmap)
        return scene

    def show_next_slide(self):
        self.current_slide_index = (self.current_slide_index + 1) % len(self.presentation.slides)
        self.show_current_slide()

if __name__ == "__main__":
    # Example usage
    pptx_file_path = "Speed of COVID-19 Vaccine Development.pptx"  # Replace with the actual path to your .pptx file
    presentation_widget = PresentationScreen(pptx_file_path)
    presentation_widget.show()  # Show the widget
