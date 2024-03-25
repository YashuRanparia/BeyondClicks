import os
import sys
import tkinter as tk
from tkinter import filedialog

import win32com.client
from pptx import Presentation
from pptx.util import Inches, Pt
from PyQt5 import QtWidgets
from PyQt5.QtCore import Qt, pyqtSlot
from PyQt5.QtGui import QPixmap
from PyQt5.QtWidgets import (QAction, QApplication, QDesktopWidget,
                             QFileDialog, QGraphicsPixmapItem, QGraphicsScene,
                             QGraphicsView, QLabel, QListWidget,
                             QListWidgetItem, QMainWindow, QMenuBar,
                             QPushButton, QStackedWidget, QTextEdit,
                             QVBoxLayout, QWidget)

from filesView import FilesView
from VirtualWrittingBoard import VideoThread as vt
from VirtualWrittingBoard import VWBView as vwb

#Static sizes
sizes = {}
sizes["smartppt"] = (1000,40)


#Defining my own window
class MyWindow(QWidget):
    def __init__(self):
        super().__init__()
        self.initUI()
        
    def initUI(self):
        # Create a QListView widget
        list_view = QListWidget(self)
        list_view.setMinimumSize(1000,400)
        list_view.move(350,80)
        list_view.itemClicked.connect(self.onItemClick)

        #Temperarory push button
        button = QPushButton(self)
        button.setText("Find files!")
        button.clicked.connect(self.find_files)
        button.move(0,50)


        # #making the text field for the single file item
        text_widget = QTextEdit(self)
        text_widget.setReadOnly(True)
        text_widget.move(0,80)
        text_widget.setMinimumSize(300,300)
        text_widget.setViewportMargins(50,20,50,20)


        #Set the Layout
        layout = QVBoxLayout()
        layout.addWidget(button)
        layout.addWidget(text_widget)
        layout.addWidget(list_view)
        
        self.setLayout(layout)

        #OnCLick button event: Select a folder
    def find_files(self):
        folder_path = QFileDialog.getExistingDirectory(self, "Select a folder")
        if folder_path:
            pptx_files = find_pptx_files(folder_path)
            if pptx_files:
                self.text_widget.append("\n\nPPTX Files Found:")
                for pptx_file in pptx_files:
                    self.list_view.addItem(pptx_file)
            else:
                self.text_widget.append("No .pptx files found in the selected folder.")
        else:
            self.text_widget.append("No folder selected.")

    def onItemClick(self,item):
        print("Clicked on: " + item.text())
        self.screen2 = QWidget(self)
        layout_2 = QVBoxLayout()
        text = QTextEdit()
        text.setReadOnly(True)
        text.setFixedSize(900,600)
        text.append(item.text())
        layout_2.addWidget(text)
        self.screen2.setLayout(layout_2)
        self.setCentralWidget(self.screen2)

    # def update(self):
    #     self.label.adjustSize()
    


class PresentationScreen(QWidget):
    
    def __init__(self):
        super().__init__()
        pptx_file_path = "Speed of COVID-19 Vaccine Development.pptx"
        self.presentation = Presentation(pptx_file_path)
        self.current_slide_index = 0
        print(len(self.presentation.slides))
        self.setup_ui()

    def setup_ui(self):
        layout = QVBoxLayout()

        # Display the current slide's content
        self.slide_view = QGraphicsView()
        layout.addWidget(self.slide_view)

        # Next Slide button
        next_button = QPushButton("Next Slide")
        next_button.clicked.connect(self.show_next_slide)
        layout.addWidget(next_button)

        self.setLayout(layout)
        self.show_current_slide()

    def show_current_slide(self):
        current_slide = self.presentation.slides[self.current_slide_index]
        slide_image = self.get_slide_image(current_slide)
        self.slide_view.setScene(slide_image)

    def render_slide(self, slide):
        scene = QGraphicsScene()
        pixmap_item = QGraphicsPixmapItem(self.get_slide_image(slide))
        scene.addItem(pixmap_item)
        return scene

    def get_slide_image(self, slide):
        # Calculate slide dimensions (adjust as needed)
        slide_width = Inches(6.5)
        slide_height = Inches(4.5)

        # Extract the slide image (first shape with an image)
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Shape type for images
                image_stream = shape.image.blob
                pixmap = QPixmap()
                pixmap.loadFromData(image_stream)
                pixmap = pixmap.scaled(slide_width, slide_height, aspectMode=Qt.KeepAspectRatio, transformMode=Qt.SmoothTransformation)
                scene = QGraphicsScene()
                scene.addPixmap(pixmap)
                return scene

    def show_next_slide(self):
        self.current_slide_index = (self.current_slide_index + 1) % len(self.presentation.slides)
        self.show_current_slide()

#Finding the .pptx files from the directories
def find_pptx_files(directory):
    pptx_files = []
    for root, _, files in os.walk(directory):
        for filename in files:
            if filename.lower().endswith(".pptx"):
                pptx_files.append(os.path.join(root, filename))
    return pptx_files


# Entry Point: Main Window
class MainApp(QMainWindow):

    def __init__(self):
        super().__init__()
        screen_geometry = QDesktopWidget().screenGeometry()
        self.setGeometry(screen_geometry)
        self.showMaximized()
        self.setWindowTitle("SGP4")

        # Create a stacked widget to manage slides
        self.stacked_widget = QStackedWidget()
        self.setCentralWidget(self.stacked_widget)

        self.page1 = MyWindow()
        # self.pres = PresentationScreen()
        self.fileView = FilesView("C:\\Users\\yashu\\Desktop\\SGP_4\\SGP4\\Test")

        self.initUI()
        pass


    def initUI(self):
        self.createMenubar()
        self.stacked_widget.addWidget(self.fileView)
        pass

    #Menubar
    def createMenubar(self):

        #Menubar:- File
        menubar = self.menuBar()
        file_menu = menubar.addMenu("File")
        open_action = QAction("Open",self)
        file_menu.addAction(open_action)
        
        #Menubar:- VWB
        presentation_menu = menubar.addMenu("VWB")
        presentation_action = QAction("Write Virtually",self)
        presentation_action.triggered.connect(self.openVWB)
        presentation_menu.addAction(presentation_action)
        
        #Menubar:- Help
        menubar = self.menuBar()
        help_menu = menubar.addMenu("Help")
        about_action = QAction("About",self)
        file_menu.addAction(about_action)

        # Connect actions to functions
        open_action.triggered.connect(self.selectDir)
        about_action.triggered.connect(self.show_about)
        presentation_action.triggered.connect(self.show_presentation)


    def selectDir(self):
        options = QFileDialog.Options()
        folder_path = QFileDialog.getExistingDirectory(self, "Select Folder", options=options)
        print(folder_path)
        self.fileView.updateUI(folder_path)
        pass

    #Show About
    def show_about(self):
        self.text_widget.append("PLease click on the file to select the file and then select the file to give presentation.")

    def openVWB(self):
        # self.showMinimized()

        vi = vwb()

        self.stacked_widget.addWidget(vi)
        self.stacked_widget.setCurrentWidget(vi)
        vi.start()

        # self.show(obj)
        pass

    #Presentation
    def show_presentation(self):
        pass

    #Slides
    def switch_to_slide(self, slide_number):
        # Show the specified slide
        self.stacked_widget.setCurrentIndex(slide_number - 1)

def onclick():
    print("Finding files.......")



if __name__ == "__main__":
    app = QApplication(sys.argv)
    dirPath = 'C:\\Users\\yashu\\Desktop\\SGP_4\\SGP4\\Test'
    win = MainApp()
    win.show()
    sys.exit(app.exec_())
